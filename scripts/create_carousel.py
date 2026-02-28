"""
Creates the OpenComposites LinkedIn/Instagram carousel as a .pptx file.
10 slides, 1080x1080px (square), brand colours.

Usage:
    python scripts/create_carousel.py
Output:
    marketing/opencomposites-carousel.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
import pptx.util as util

# ── Brand Colours ────────────────────────────────────────────────────────────
SLATE_DARK  = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A — background
SLATE_MID   = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B — card bg
ORANGE      = RGBColor(0xF9, 0x73, 0x16)   # #F97316 — accent
RED         = RGBColor(0xE8, 0x4C, 0x3D)   # #E84C3D — primary
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0x94, 0xA3, 0xB8)   # #94A3B8 — subtext

# ── Slide size: 1080x1080 at 96dpi = 11.25in x 11.25in ─────────────────────
SLIDE_W = Inches(11.25)
SLIDE_H = Inches(11.25)

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    """Add a completely blank slide."""
    blank_layout = prs.slide_layouts[6]   # index 6 = blank
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, colour=SLATE_DARK):
    """Solid background fill for the whole slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour

def add_rect(slide, left, top, width, height, colour):
    """Add a filled rectangle (no border)."""
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()   # no border
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=32, bold=False, colour=WHITE,
             align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = "Calibri"   # closest to Inter available in pptx
    return txBox

def add_label(slide, text, left, top, width,
              font_size=11, colour=LIGHT_GREY, align=PP_ALIGN.LEFT):
    """Small ALL-CAPS label."""
    return add_text(slide, text.upper(), left, top, width, Inches(0.4),
                    font_size=font_size, bold=True, colour=colour, align=align)

def add_pill(slide, text, left, top):
    """Orange pill badge."""
    w, h = Inches(2.8), Inches(0.45)
    add_rect(slide, left, top, w, h, ORANGE)
    add_text(slide, text, left + Inches(0.15), top + Inches(0.04),
             w - Inches(0.3), h, font_size=13, bold=True,
             colour=WHITE, align=PP_ALIGN.CENTER)

def add_footer(slide, text="by Addcomposites · addcomposites.com",
               accent=False):
    """Consistent footer across slides."""
    add_text(slide, text,
             Inches(0.5), SLIDE_H - Inches(0.7),
             Inches(10.25), Inches(0.55),
             font_size=11, colour=LIGHT_GREY,
             align=PP_ALIGN.LEFT)
    # Orange accent line at bottom
    add_rect(slide,
             0, SLIDE_H - Inches(0.08),
             SLIDE_W, Inches(0.08),
             ORANGE)

def add_slide_number(slide, n):
    """Small slide number top-right."""
    add_text(slide, f"{n}/10",
             SLIDE_W - Inches(1.2), Inches(0.35),
             Inches(0.9), Inches(0.4),
             font_size=13, colour=LIGHT_GREY,
             align=PP_ALIGN.RIGHT)


# ── Slide Builders ───────────────────────────────────────────────────────────

def slide_1_hook(prs):
    """Hook / Cover."""
    s = blank_slide(prs)
    fill_bg(s)

    # Top accent bar
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), ORANGE)

    # Main headline
    add_text(s,
             "99% of people interested in composites have no access to the tools that matter.",
             Inches(0.7), Inches(1.8), Inches(9.85), Inches(4.5),
             font_size=52, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    # Sub
    add_text(s, "We fixed that. For free.",
             Inches(0.7), Inches(6.5), Inches(9.85), Inches(1.2),
             font_size=32, bold=False, colour=ORANGE, align=PP_ALIGN.CENTER)

    add_footer(s, "OpenComposites · by Addcomposites")
    return s


def slide_2_problem(prs):
    """The problem."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), RED)

    add_text(s, "The software that answers composites questions costs:",
             Inches(0.7), Inches(0.8), Inches(9.85), Inches(1.1),
             font_size=26, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

    tools = [
        ("Fibersim", "$15,000+ / year"),
        ("CATIA Composites", "$50,000+ / year"),
        ("HyperFiber", "$30,000+ / year"),
    ]
    y = Inches(2.2)
    for tool, price in tools:
        # card
        add_rect(s, Inches(1.0), y, Inches(9.25), Inches(1.5), SLATE_MID)
        add_text(s, tool,
                 Inches(1.4), y + Inches(0.18), Inches(5.0), Inches(0.75),
                 font_size=24, bold=True, colour=WHITE)
        add_text(s, price,
                 Inches(6.5), y + Inches(0.18), Inches(3.2), Inches(0.75),
                 font_size=24, bold=True, colour=RED, align=PP_ALIGN.RIGHT)
        y += Inches(1.7)

    add_text(s, "None of these are accessible to makers, drone startups, or junior engineers.",
             Inches(0.7), Inches(7.3), Inches(9.85), Inches(1.2),
             font_size=20, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_slide_number(s, 2)
    add_footer(s)
    return s


def slide_3_solution(prs):
    """Introduce the solution."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), ORANGE)

    add_text(s, "So we built OpenComposites.",
             Inches(0.7), Inches(1.5), Inches(9.85), Inches(1.6),
             font_size=56, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    badges = ["Free", "Open Source", "No account needed"]
    x = Inches(1.2)
    for badge in badges:
        add_pill(s, badge, x, Inches(3.6))
        x += Inches(3.1)

    add_text(s,
             "Built by the team at Addcomposites — the company that makes AFP robots for industrial composites manufacturing.",
             Inches(0.7), Inches(4.6), Inches(9.85), Inches(1.5),
             font_size=22, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_text(s, "opencomposites.addcomposites.com",
             Inches(0.7), Inches(6.4), Inches(9.85), Inches(0.9),
             font_size=26, bold=True, colour=ORANGE, align=PP_ALIGN.CENTER)

    add_slide_number(s, 3)
    add_footer(s)
    return s


def slide_4_ai(prs):
    """Feature: AI Part Analysis."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), ORANGE)

    add_label(s, "Feature 1", Inches(0.7), Inches(0.5), Inches(4.0))
    add_text(s, "Describe your part.\nGet a full design plan.",
             Inches(0.7), Inches(0.9), Inches(9.85), Inches(2.0),
             font_size=42, bold=True, colour=WHITE)

    # Input mock
    add_rect(s, Inches(0.7), Inches(3.2), Inches(9.85), Inches(1.2), SLATE_MID)
    add_text(s, '"200mm drone arm, 50N lateral load, carbon fibre, lightest possible"',
             Inches(0.9), Inches(3.3), Inches(9.45), Inches(1.0),
             font_size=18, colour=LIGHT_GREY)

    # Arrow
    add_text(s, "→", Inches(5.0), Inches(4.6), Inches(1.2), Inches(0.7),
             font_size=36, bold=True, colour=ORANGE, align=PP_ALIGN.CENTER)

    # Output bullets
    outputs = [
        "Material selection + reasoning",
        "Fibre orientation & stacking sequence",
        "Manufacturing process recommendation",
        "Per-part cost estimate",
    ]
    y = Inches(5.5)
    for o in outputs:
        add_rect(s, Inches(0.7), y, Inches(0.35), Inches(0.35), ORANGE)
        add_text(s, o, Inches(1.2), y - Inches(0.04), Inches(8.7), Inches(0.5),
                 font_size=20, colour=WHITE)
        y += Inches(0.65)

    add_text(s, "Powered by Claude AI · BYOK (bring your own Anthropic key)",
             Inches(0.7), SLIDE_H - Inches(1.4), Inches(9.85), Inches(0.55),
             font_size=14, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_slide_number(s, 4)
    add_footer(s)
    return s


def slide_5_knowledge(prs):
    """Feature: Knowledge Base."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), ORANGE)

    add_label(s, "Feature 2", Inches(0.7), Inches(0.5), Inches(4.0))
    add_text(s, "56 plain-language articles.\nCC BY 4.0. Forever free.",
             Inches(0.7), Inches(0.9), Inches(9.85), Inches(2.2),
             font_size=42, bold=True, colour=WHITE)

    categories = [
        ("Fibre & Resin Science", "01 Fundamentals"),
        ("Design Rules", "02 Stacking, drop-offs, splices"),
        ("Manufacturing", "03 Wet layup → AFP"),
        ("Structural Analysis", "04 CLT, failure criteria"),
        ("CATIA Workflows", "05 Zone, ply, flat pattern"),
        ("Cost & Applications", "08–10 Real-world cases"),
    ]

    x_positions = [Inches(0.6), Inches(3.85), Inches(7.1)]
    y_row1 = Inches(3.4)
    y_row2 = Inches(5.8)

    for i, (title, sub) in enumerate(categories):
        col = i % 3
        row = i // 3
        x = x_positions[col]
        y = y_row1 if row == 0 else y_row2

        add_rect(s, x, y, Inches(2.85), Inches(1.9), SLATE_MID)
        add_text(s, title, x + Inches(0.15), y + Inches(0.18),
                 Inches(2.55), Inches(0.85),
                 font_size=16, bold=True, colour=WHITE)
        add_text(s, sub, x + Inches(0.15), y + Inches(1.05),
                 Inches(2.55), Inches(0.7),
                 font_size=13, colour=LIGHT_GREY)

    add_text(s, "Structured for LLM retrieval — Claude, ChatGPT, Gemini, local models can all search it.",
             Inches(0.7), SLIDE_H - Inches(1.4), Inches(9.85), Inches(0.55),
             font_size=14, colour=ORANGE, align=PP_ALIGN.CENTER)

    add_slide_number(s, 5)
    add_footer(s)
    return s


def slide_6_calculators(prs):
    """Feature: Engineering Calculators."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), ORANGE)

    add_label(s, "Feature 3", Inches(0.7), Inches(0.5), Inches(4.0))
    add_text(s, "Not just AI — real\nengineering tools.",
             Inches(0.7), Inches(0.9), Inches(9.85), Inches(2.2),
             font_size=44, bold=True, colour=WHITE)

    calcs = [
        ("CLT Calculator", "ABD matrices, effective moduli,\nfailure criteria (Tsai-Wu, Hashin)"),
        ("Sandwich Panel", "Core + face sheet design\nand optimisation"),
        ("Bolted Joint", "Bearing & bypass analysis,\njoint sizing"),
        ("Cost Estimator", "Material, labour, tooling\nper-part cost breakdown"),
    ]

    x_positions = [Inches(0.6), Inches(5.55)]
    y_positions = [Inches(3.3), Inches(6.0)]

    for i, (title, desc) in enumerate(calcs):
        col = i % 2
        row = i // 2
        x = x_positions[col]
        y = y_positions[row]

        add_rect(s, x, y, Inches(4.55), Inches(2.2), SLATE_MID)
        # Orange top stripe
        add_rect(s, x, y, Inches(4.55), Inches(0.12), ORANGE)
        add_text(s, title, x + Inches(0.2), y + Inches(0.25),
                 Inches(4.15), Inches(0.7),
                 font_size=20, bold=True, colour=WHITE)
        add_text(s, desc, x + Inches(0.2), y + Inches(0.95),
                 Inches(4.15), Inches(1.1),
                 font_size=15, colour=LIGHT_GREY)

    add_slide_number(s, 6)
    add_footer(s)
    return s


def slide_7_personas(prs):
    """Who it's for."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), ORANGE)

    add_text(s, "Built for the 99% who don't have access.",
             Inches(0.7), Inches(0.7), Inches(9.85), Inches(1.3),
             font_size=38, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    personas = [
        ("The Maker",
         '"I want to build a carbon fibre part for my car. I have no idea which fibres to buy."'),
        ("The Drone Engineer",
         '"Designing an airframe at a startup — no composites person on the team, tight budget."'),
        ("The Junior Engineer",
         '"My company can\'t afford Fibersim. I need composites knowledge somewhere reliable."'),
    ]

    x = Inches(0.55)
    for title, quote in personas:
        add_rect(s, x, Inches(2.3), Inches(3.2), Inches(6.8), SLATE_MID)
        add_rect(s, x, Inches(2.3), Inches(3.2), Inches(0.12), ORANGE)
        add_text(s, title, x + Inches(0.2), Inches(2.6),
                 Inches(2.8), Inches(0.7),
                 font_size=20, bold=True, colour=ORANGE)
        add_text(s, quote, x + Inches(0.2), Inches(3.45),
                 Inches(2.8), Inches(5.0),
                 font_size=16, colour=WHITE)
        x += Inches(3.4)

    add_slide_number(s, 7)
    add_footer(s)
    return s


def slide_8_opensource(prs):
    """Open source angle."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), ORANGE)

    add_text(s, "Open source.\nCitable. Forkable.",
             Inches(0.7), Inches(0.9), Inches(9.85), Inches(2.5),
             font_size=50, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    points = [
        ("CC BY 4.0",     "Use in your thesis, product docs, or startup's knowledge base"),
        ("Fork it",       "Adapt it for your domain — FRP civil, sporting goods, marine"),
        ("Contribute",    "Add an article, fix a mistake, get credited as a contributor"),
        ("LLM-ready",     "RAG-structured — any AI assistant can search it reliably"),
    ]

    y = Inches(3.8)
    for label, desc in points:
        add_rect(s, Inches(0.7), y + Inches(0.08), Inches(2.2), Inches(0.72), ORANGE)
        add_text(s, label, Inches(0.8), y + Inches(0.08),
                 Inches(2.0), Inches(0.72),
                 font_size=17, bold=True, colour=SLATE_DARK, align=PP_ALIGN.CENTER)
        add_text(s, desc, Inches(3.1), y + Inches(0.1),
                 Inches(7.5), Inches(0.68),
                 font_size=17, colour=WHITE)
        y += Inches(1.1)

    add_slide_number(s, 8)
    add_footer(s)
    return s


def slide_9_credibility(prs):
    """Credibility / who built it."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), ORANGE)

    add_text(s, "Built by people who make\ncomposites machines for a living.",
             Inches(0.7), Inches(0.9), Inches(9.85), Inches(2.5),
             font_size=44, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(0.7), Inches(3.8), Inches(9.85), Inches(3.5), SLATE_MID)

    add_text(s, "Addcomposites",
             Inches(1.1), Inches(4.0), Inches(9.0), Inches(0.9),
             font_size=28, bold=True, colour=ORANGE)

    add_text(s,
             "We build AFP (Automated Fibre Placement) robots used in industrial composites manufacturing. "
             "We've spent years watching engineers struggle with composites questions we could answer in 5 minutes — "
             "but the right tools cost $15,000+ a year. This is our response.",
             Inches(1.1), Inches(4.95), Inches(9.0), Inches(2.1),
             font_size=19, colour=WHITE)

    add_text(s, "addcomposites.com — AFP-XS · AFP-X · AddPath software",
             Inches(0.7), Inches(7.75), Inches(9.85), Inches(0.55),
             font_size=15, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_slide_number(s, 9)
    add_footer(s)
    return s


def slide_10_cta(prs):
    """Call to action."""
    s = blank_slide(prs)
    fill_bg(s, ORANGE)   # orange background for the CTA slide

    # Overlay slate for contrast
    add_rect(s, 0, 0, SLIDE_W, Inches(1.8), SLATE_DARK)
    add_text(s, "Try it free.",
             Inches(0.7), Inches(0.3), Inches(9.85), Inches(1.3),
             font_size=72, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    add_rect(s, 0, Inches(1.8), SLIDE_W, SLIDE_H - Inches(1.8), SLATE_DARK)

    add_text(s, "opencomposites.addcomposites.com",
             Inches(0.7), Inches(2.3), Inches(9.85), Inches(1.2),
             font_size=30, bold=True, colour=ORANGE, align=PP_ALIGN.CENTER)

    features = [
        "No account required",
        "No credit card",
        "No paywall — ever",
        "Open source (CC BY 4.0)",
    ]
    y = Inches(4.0)
    for f in features:
        add_text(s, f"✓  {f}",
                 Inches(2.5), y, Inches(6.25), Inches(0.65),
                 font_size=22, colour=WHITE, align=PP_ALIGN.LEFT)
        y += Inches(0.72)

    add_text(s, "github.com/Addcomposites-github/composites-design-guide",
             Inches(0.7), Inches(7.8), Inches(9.85), Inches(0.6),
             font_size=15, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_rect(s, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ORANGE)
    add_text(s, "OpenComposites · by Addcomposites",
             Inches(0.7), SLIDE_H - Inches(0.7), Inches(9.85), Inches(0.55),
             font_size=11, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)
    return s


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    out_dir = os.path.join(os.path.dirname(__file__), "..", "marketing")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "opencomposites-carousel.pptx")

    prs = new_prs()
    slide_1_hook(prs)
    slide_2_problem(prs)
    slide_3_solution(prs)
    slide_4_ai(prs)
    slide_5_knowledge(prs)
    slide_6_calculators(prs)
    slide_7_personas(prs)
    slide_8_opensource(prs)
    slide_9_credibility(prs)
    slide_10_cta(prs)

    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
